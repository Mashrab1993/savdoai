"""
╔══════════════════════════════════════════════════════════════════════════╗
║  SAVDOAI — HUJJAT O'QISH MEGA v3.0 (100,000+ SAHIFA)                   ║
║                                                                          ║
║  DISK-BASED ARXITEKTURA:                                                ║
║  - Sahifa matni DISKDA saqlanadi (xotira emas)                          ║
║  - Faqat INDEKS xotirada (sahifa boshlanishi + kalit so'zlar)          ║
║  - Lazy loading — sahifa faqat so'ralganda diskdan o'qiladi            ║
║  - 100,000 sahifali kitob = ~10MB indeks + disk cache                  ║
║                                                                          ║
║  32 FORMAT: PDF, DOCX, XLSX, EPUB, PPTX, FB2, RTF, HTML, JSON, MD...  ║
╚══════════════════════════════════════════════════════════════════════════╝
"""
from __future__ import annotations
import io, json, logging, os, re, tempfile, time, hashlib, sqlite3
from typing import Optional
from collections import Counter
log = logging.getLogger(__name__)

# Disk cache papkasi
_CACHE_DIR = os.path.join(tempfile.gettempdir(), "savdoai_hujjat")
os.makedirs(_CACHE_DIR, exist_ok=True)
_CACHE_MAX_AGE = 3600 * 6  # 6 soat


def cache_tozalash():
    """Eski cache fayllarni o'chirish (6 soatdan eski)."""
    try:
        now = time.time()
        for f in os.listdir(_CACHE_DIR):
            fp = os.path.join(_CACHE_DIR, f)
            if os.path.isfile(fp) and now - os.path.getmtime(fp) > _CACHE_MAX_AGE:
                os.unlink(fp)
    except Exception:
        pass


# ═══════════════════════════════════════════════════════════════
#  DISK CACHE — sahifa matnini diskda saqlash
# ═══════════════════════════════════════════════════════════════

class DiskCache:
    """SQLite-based sahifa cache — 100K+ sahifa, ~0 xotira."""
    
    def __init__(self, fayl_id: str):
        self.db_path = os.path.join(_CACHE_DIR, f"{fayl_id}.db")
        self.conn = sqlite3.connect(self.db_path)
        self.conn.execute("CREATE TABLE IF NOT EXISTS sahifalar (num INTEGER PRIMARY KEY, matn TEXT)")
        self.conn.execute("CREATE TABLE IF NOT EXISTS indeks (num INTEGER PRIMARY KEY, boshi TEXT, sozlar TEXT)")
        # FTS5 — tezkor full-text search
        try:
            self.conn.execute("CREATE VIRTUAL TABLE IF NOT EXISTS fts USING fts5(num, matn, boshi, tokenize='unicode61')")
        except Exception:
            pass  # FTS5 mavjud bo'lmasa oddiy LIKE ishlaydi
        self.conn.execute("PRAGMA journal_mode=WAL")
        self.conn.execute("PRAGMA synchronous=NORMAL")
    
    def sahifa_yoz(self, num: int, matn: str):
        self.conn.execute("INSERT OR REPLACE INTO sahifalar VALUES (?, ?)", (num, matn))
        try:
            self.conn.execute("INSERT OR REPLACE INTO fts(num, matn, boshi) VALUES (?, ?, ?)",
                              (str(num), matn[:500], matn[:150]))
        except Exception:
            pass
    
    def indeks_yoz(self, num: int, boshi: str, sozlar: list):
        self.conn.execute("INSERT OR REPLACE INTO indeks VALUES (?, ?, ?)",
                          (num, boshi, json.dumps(sozlar, ensure_ascii=False)))
    
    def commit(self):
        self.conn.commit()
    
    def sahifa_oqi(self, num: int) -> str:
        row = self.conn.execute("SELECT matn FROM sahifalar WHERE num=?", (num,)).fetchone()
        return row[0] if row else ""
    
    def indeks_oqi(self, num: int) -> dict:
        row = self.conn.execute("SELECT boshi, sozlar FROM indeks WHERE num=?", (num,)).fetchone()
        if row:
            try: sozlar = json.loads(row[1])
            except: sozlar = []
            return {"boshi": row[0], "sozlar": sozlar}
        return {"boshi": "", "sozlar": []}
    
    def barcha_indeks(self) -> dict:
        rows = self.conn.execute("SELECT num, boshi, sozlar FROM indeks ORDER BY num").fetchall()
        result = {}
        for num, boshi, sozlar_json in rows:
            try: sozlar = json.loads(sozlar_json)
            except: sozlar = []
            result[num] = {"boshi": boshi, "sozlar": sozlar}
        return result
    
    def izla(self, kalit_sozlar: list, limit: int = 50) -> list:
        """FTS5 orqali tezkor izlash — 100K da ham 1ms."""
        if not kalit_sozlar:
            return []
        topilgan = []
        # FTS5 izlash (eng tez)
        try:
            query = " OR ".join(kalit_sozlar)
            rows = self.conn.execute(
                "SELECT num FROM fts WHERE fts MATCH ? ORDER BY rank LIMIT ?",
                (query, limit)).fetchall()
            topilgan = [int(r[0]) for r in rows]
            if topilgan:
                return topilgan
        except Exception:
            pass
        # Fallback: indeks LIKE
        for k in kalit_sozlar:
            rows = self.conn.execute(
                "SELECT num FROM indeks WHERE boshi LIKE ? OR sozlar LIKE ? ORDER BY num LIMIT ?",
                (f"%{k}%", f"%{k}%", limit)).fetchall()
            for (num,) in rows:
                if num not in topilgan:
                    topilgan.append(num)
        if topilgan:
            return topilgan[:limit]
        # Fallback: to'liq matn LIKE
        if len(kalit_sozlar) == 1:
            rows = self.conn.execute(
                "SELECT num FROM sahifalar WHERE matn LIKE ? ORDER BY num LIMIT ?",
                (f"%{kalit_sozlar[0]}%", limit)).fetchall()
            return [num for (num,) in rows]
        return []
    
    def close(self):
        self.conn.close()
    
    def __del__(self):
        try: self.conn.close()
        except: pass


def _fayl_id(data: bytes) -> str:
    return hashlib.md5(data[:10000]).hexdigest()[:16]

def _kalit_sozlar(matn: str, limit: int = 20) -> list:
    if not matn: return []
    sozlar = re.findall(r'\b[A-Za-zА-Яа-яЎўҚқҒғҲҳ]{3,}\b', matn.lower())
    c = Counter(sozlar)
    stop = {"bir","va","bilan","uchun","dan","ning","the","and","for","is","this","that","with",
            "это","для","при","что","как","или","не","по","на","из","har","shu","bu","yoki","agar","ham","kerak"}
    return [w for w, _ in c.most_common(limit + len(stop)) if w not in stop][:limit]

def _sarlavha_ekanmi(qator: str) -> bool:
    q = qator.strip()
    if not q or len(q) < 3 or len(q) > 100: return False
    if re.match(r'^(§|\d+[\.\)]\s|Глава|Bob|CHAPTER|Chapter|Mavzu|MAVZU|Раздел|PART|Part)', q): return True
    if q.isupper() and len(q) > 5: return True
    if not q.endswith('.') and len(q) < 60 and q[0].isupper(): return True
    return False


# ═══════════════════════════════════════════════════════════════
#  PDF — 100,000+ SAHIFA (DISK-BASED)
# ═══════════════════════════════════════════════════════════════

def pdf_oqi(data: bytes) -> dict:
    try:
        import pdfplumber
    except ImportError:
        return {"xato": "pdfplumber kerak", "sahifalar_soni": 0}
    
    fid = _fayl_id(data)
    cache = DiskCache(fid)
    result = {"sahifalar_soni": 0, "mundarija": [], "jadvallar": [], "metadata": {},
              "_cache": cache, "_fid": fid, "fayl_data": data, "umumiy_matn": ""}
    
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            jami = len(pdf.pages)
            result["sahifalar_soni"] = jami
            result["metadata"] = pdf.metadata or {}
            umumiy = []
            
            for i, page in enumerate(pdf.pages, 1):
                matn = (page.extract_text() or "").strip()
                
                # Diskga yozish
                cache.sahifa_yoz(i, matn)
                boshi = matn[:150].replace("\n", " ").strip() if matn else ""
                cache.indeks_yoz(i, boshi, _kalit_sozlar(matn))
                
                # Mundarija
                if matn:
                    for qator in matn.split("\n")[:5]:
                        if _sarlavha_ekanmi(qator.strip()):
                            result["mundarija"].append({"sahifa": i, "sarlavha": qator.strip()[:80]})
                            break
                
                # Jadvallar (birinchi 30 sahifadan)
                if i <= 30:
                    try:
                        for t_idx, table in enumerate(page.extract_tables()):
                            if table and len(table) > 1:
                                result["jadvallar"].append({"sahifa": i, "jadval_raqam": t_idx+1,
                                    "sarlavha": table[0], "qatorlar": table[1:6], "qator_soni": len(table)})
                    except: pass
                
                # Umumiy — birinchi 5 + oxirgi 3
                if i <= 5 or i > jami - 3:
                    umumiy.append(matn[:300])
                
                # Progress
                if jami > 500 and i % 500 == 0:
                    log.info("PDF: %d/%d (%.0f%%)...", i, jami, i/jami*100)
                    cache.commit()
            
            cache.commit()
            result["umumiy_matn"] = "\n\n".join(umumiy)
            log.info("PDF: %d sahifa, %d mundarija, %d jadval", jami, len(result["mundarija"]), len(result["jadvallar"]))
    except Exception as e:
        log.error("PDF: %s", e); result["xato"] = str(e)
    return result


# ═══════════════════════════════════════════════════════════════
#  WORD — 100,000+ PARAGRAF
# ═══════════════════════════════════════════════════════════════

def docx_oqi(data: bytes) -> dict:
    try:
        from docx import Document
    except ImportError:
        return {"xato": "python-docx kerak"}
    fid = _fayl_id(data)
    cache = DiskCache(fid)
    result = {"paragraflar_soni": 0, "sarlavhalar": [], "jadvallar": [],
              "mundarija": [], "sahifalar_soni": 0, "_cache": cache, "_fid": fid, "umumiy_matn": ""}
    try:
        doc = Document(io.BytesIO(data))
        matn_q = []; sahifa_num = 1; buf = []; buf_len = 0
        for i, para in enumerate(doc.paragraphs, 1):
            matn = para.text.strip()
            if not matn: continue
            stil = para.style.name if para.style else ""
            if "Heading" in stil:
                daraja = 1
                try: daraja = int(re.search(r'\d', stil).group())
                except: pass
                result["sarlavhalar"].append({"daraja": daraja, "matn": matn})
                result["mundarija"].append({"sahifa": sahifa_num, "sarlavha": matn[:80]})
            buf.append(matn); buf_len += len(matn)
            if len(buf) >= 40 or buf_len >= 2000:
                sahifa_matn = "\n".join(buf)
                cache.sahifa_yoz(sahifa_num, sahifa_matn)
                cache.indeks_yoz(sahifa_num, sahifa_matn[:150].replace("\n"," "), _kalit_sozlar(sahifa_matn))
                matn_q.append(sahifa_matn[:300])
                sahifa_num += 1; buf = []; buf_len = 0
        if buf:
            sm = "\n".join(buf)
            cache.sahifa_yoz(sahifa_num, sm)
            cache.indeks_yoz(sahifa_num, sm[:150].replace("\n"," "), _kalit_sozlar(sm))
        cache.commit()
        result["paragraflar_soni"] = i if 'i' in dir() else 0
        result["sahifalar_soni"] = sahifa_num
        for t_idx, table in enumerate(doc.tables[:20]):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows[:50]]
            if rows: result["jadvallar"].append({"jadval_raqam": t_idx+1, "sarlavha": rows[0], "qatorlar": rows[1:6], "qator_soni": len(rows)})
        result["umumiy_matn"] = "\n".join(matn_q[:20])
    except Exception as e:
        log.error("DOCX: %s", e); result["xato"] = str(e)
    return result


# ═══════════════════════════════════════════════════════════════
#  EXCEL — 100,000+ QATOR
# ═══════════════════════════════════════════════════════════════

def xlsx_oqi(data: bytes) -> dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"xato": "openpyxl kerak"}
    fid = _fayl_id(data)
    cache = DiskCache(fid)
    result = {"sheetlar_soni": 0, "sheetlar": [], "sahifalar_soni": 0, "_cache": cache, "_fid": fid, "umumiy_matn": ""}
    try:
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        result["sheetlar_soni"] = len(wb.sheetnames); sahifa_num = 1; umumiy = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]; rows = []; raqamlar = []
            for idx, row in enumerate(ws.iter_rows(values_only=True)):
                rows.append(row)
                for c in row:
                    if isinstance(c, (int, float)) and c != 0: raqamlar.append(c)
                if idx >= 100000: break
            if not rows: continue
            sarlavha = [str(c or "") for c in rows[0]]
            for cs in range(0, len(rows), 100):
                chunk = rows[cs:cs+100]
                sm = "\n".join(" | ".join(str(c or "") for c in r) for r in chunk)
                cache.sahifa_yoz(sahifa_num, sm)
                cache.indeks_yoz(sahifa_num, sm[:150], _kalit_sozlar(sm))
                sahifa_num += 1
            result["sheetlar"].append({"nom": sheet_name, "qator_soni": len(rows), "ustun_soni": len(sarlavha),
                "sarlavha": sarlavha[:20],
                "statistika": {"jami_raqamlar": len(raqamlar),
                    "eng_katta": max(raqamlar) if raqamlar else 0,
                    "eng_kichik": min(r for r in raqamlar if r > 0) if [r for r in raqamlar if r > 0] else 0,
                    "ortacha": round(sum(raqamlar)/len(raqamlar), 2) if raqamlar else 0,
                    "jami": round(sum(raqamlar), 2) if raqamlar else 0}})
            for r in rows[:20]: umumiy.append(" | ".join(str(c or "") for c in r))
        cache.commit()
        result["sahifalar_soni"] = sahifa_num - 1; result["umumiy_matn"] = "\n".join(umumiy)
        wb.close()
    except Exception as e:
        log.error("XLSX: %s", e); result["xato"] = str(e)
    return result


# ═══════════════════════════════════════════════════════════════
#  EPUB, PPTX, RTF, HTML, FB2, JSON, MD, ODT
# ═══════════════════════════════════════════════════════════════

def epub_oqi(data: bytes) -> dict:
    try:
        import ebooklib; from ebooklib import epub; from bs4 import BeautifulSoup
    except ImportError:
        return {"xato": "ebooklib + beautifulsoup4 kerak"}
    fid = _fayl_id(data); cache = DiskCache(fid)
    result = {"sahifalar_soni": 0, "mundarija": [], "jadvallar": [], "metadata": {},
              "_cache": cache, "_fid": fid, "umumiy_matn": ""}
    try:
        book = epub.read_epub(io.BytesIO(data))
        result["metadata"] = {
            "sarlavha": book.get_metadata('DC','title')[0][0] if book.get_metadata('DC','title') else "",
            "muallif": book.get_metadata('DC','creator')[0][0] if book.get_metadata('DC','creator') else "",
        }
        sahifa_num = 1; umumiy = []
        for item in book.toc[:50]:
            if hasattr(item, 'title'):
                result["mundarija"].append({"sahifa": sahifa_num, "sarlavha": str(item.title)[:80]})
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            html = item.get_content().decode('utf-8', errors='replace')
            soup = BeautifulSoup(html, 'html.parser')
            matn = soup.get_text(separator='\n', strip=True)
            if not matn.strip(): continue
            for i in range(0, len(matn), 2000):
                chunk = matn[i:i+2000]
                cache.sahifa_yoz(sahifa_num, chunk)
                cache.indeks_yoz(sahifa_num, chunk[:150].replace("\n"," "), _kalit_sozlar(chunk))
                sahifa_num += 1
            for tag in soup.find_all(['h1','h2','h3']):
                result["mundarija"].append({"sahifa": sahifa_num-1, "sarlavha": tag.get_text()[:80]})
            if len(umumiy) < 15: umumiy.append(matn[:300])
        cache.commit()
        result["sahifalar_soni"] = sahifa_num - 1; result["umumiy_matn"] = "\n\n".join(umumiy)
    except Exception as e:
        log.error("EPUB: %s", e); result["xato"] = str(e)
    return result

def pptx_oqi(data: bytes) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {"xato": "python-pptx kerak"}
    fid = _fayl_id(data); cache = DiskCache(fid)
    result = {"sahifalar_soni": 0, "mundarija": [], "jadvallar": [], "_cache": cache, "_fid": fid, "umumiy_matn": ""}
    try:
        prs = Presentation(io.BytesIO(data)); umumiy = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []; sarlavha = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt: parts.append(txt)
                        if not sarlavha and txt and len(txt) < 80: sarlavha = txt
                if shape.has_table:
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    if rows: result["jadvallar"].append({"sahifa": i, "sarlavha": rows[0], "qatorlar": rows[1:6], "qator_soni": len(rows)})
            matn = "\n".join(parts)
            cache.sahifa_yoz(i, matn)
            cache.indeks_yoz(i, matn[:150].replace("\n"," "), _kalit_sozlar(matn))
            if sarlavha: result["mundarija"].append({"sahifa": i, "sarlavha": sarlavha})
            umumiy.append(matn[:200])
        cache.commit()
        result["sahifalar_soni"] = len(prs.slides); result["umumiy_matn"] = "\n\n".join(umumiy[:20])
    except Exception as e:
        log.error("PPTX: %s", e); result["xato"] = str(e)
    return result

def _matn_cache(matn: str, tur: str, data: bytes) -> dict:
    """Universal matn → disk cache."""
    fid = _fayl_id(data); cache = DiskCache(fid)
    result = {"sahifalar_soni": 0, "mundarija": [], "jadvallar": [], "_cache": cache, "_fid": fid, "umumiy_matn": "", "metadata": {}}
    if not matn: return result
    sahifa_num = 1; buf = []; buf_len = 0
    for qator in matn.split("\n"):
        buf.append(qator); buf_len += len(qator) + 1
        if len(buf) >= 50 or buf_len >= 2000:
            sm = "\n".join(buf)
            cache.sahifa_yoz(sahifa_num, sm)
            cache.indeks_yoz(sahifa_num, sm[:150].replace("\n"," "), _kalit_sozlar(sm))
            sahifa_num += 1; buf = []; buf_len = 0
    if buf:
        sm = "\n".join(buf)
        cache.sahifa_yoz(sahifa_num, sm)
        cache.indeks_yoz(sahifa_num, sm[:150].replace("\n"," "), _kalit_sozlar(sm))
    cache.commit()
    result["sahifalar_soni"] = sahifa_num
    result["umumiy_matn"] = matn[:3000] if len(matn) <= 3000 else matn[:1500] + "\n...\n" + matn[-1000:]
    return result

def rtf_oqi(data: bytes) -> dict:
    try:
        from striprtf.striprtf import rtf_to_text
        matn = rtf_to_text(data.decode('utf-8', errors='replace'))
    except: matn = data.decode('utf-8', errors='replace')
    return _matn_cache(matn, "rtf", data)

def html_oqi(data: bytes) -> dict:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(data.decode('utf-8', errors='replace'), 'html.parser')
        matn = soup.get_text(separator='\n', strip=True)
        result = _matn_cache(matn, "html", data)
        result["mundarija"] = [{"sahifa": 0, "sarlavha": tag.get_text()[:80]} for tag in soup.find_all(['h1','h2','h3','h4'])[:50]]
        return result
    except: return _matn_cache(data.decode('utf-8', errors='replace'), "html", data)

def fb2_oqi(data: bytes) -> dict:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(data.decode('utf-8', errors='replace'), 'xml')
        parts = []; mundarija = []
        for section in soup.find_all('section'):
            title = section.find('title')
            if title: mundarija.append({"sahifa": 0, "sarlavha": title.get_text()[:80]})
            for p in section.find_all('p'): parts.append(p.get_text())
        result = _matn_cache("\n".join(parts), "fb2", data)
        result["mundarija"] = mundarija
        return result
    except: return _matn_cache(data.decode('utf-8', errors='replace'), "fb2", data)

def md_oqi(data: bytes) -> dict:
    matn = data.decode('utf-8', errors='replace')
    result = _matn_cache(matn, "md", data)
    result["mundarija"] = [{"sahifa": 0, "sarlavha": line.lstrip("#").strip()[:80]}
                            for line in matn.split("\n") if line.startswith("#") and line.lstrip("#").strip()][:50]
    return result

def json_oqi(data: bytes) -> dict:
    try: matn = json.dumps(json.loads(data), indent=2, ensure_ascii=False)
    except: matn = data.decode('utf-8', errors='replace')
    return _matn_cache(matn, "json", data)

def odt_oqi(data: bytes) -> dict:
    try:
        import zipfile; from bs4 import BeautifulSoup
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            content = zf.read('content.xml').decode('utf-8', errors='replace')
        return _matn_cache(BeautifulSoup(content, 'xml').get_text(separator='\n', strip=True), "odt", data)
    except Exception as e: return {"xato": f"ODT: {e}"}

def djvu_oqi(data: bytes) -> dict:
    return {"xato": "DJVU — rasmli format. PDF ga aylantiring.", "sahifalar_soni": 0, "tavsiya": "DJVU → PDF konvertor ishlating"}


# ═══════════════════════════════════════════════════════════════
#  UNIVERSAL ROUTER — 32 FORMAT
# ═══════════════════════════════════════════════════════════════

_FORMAT_MAP = {
    ".pdf": ("pdf", pdf_oqi), ".docx": ("docx", docx_oqi), ".doc": ("docx", docx_oqi),
    ".xlsx": ("xlsx", xlsx_oqi), ".xls": ("xlsx", xlsx_oqi),
    ".pptx": ("pptx", pptx_oqi), ".ppt": ("pptx", pptx_oqi),
    ".epub": ("epub", epub_oqi), ".fb2": ("fb2", fb2_oqi),
    ".rtf": ("rtf", rtf_oqi), ".html": ("html", html_oqi), ".htm": ("html", html_oqi),
    ".json": ("json", json_oqi), ".xml": ("xml", html_oqi),
    ".md": ("md", md_oqi), ".markdown": ("md", md_oqi),
    ".odt": ("odt", odt_oqi), ".djvu": ("djvu", djvu_oqi),
}
_TXT_EXTS = {".txt",".csv",".log",".py",".js",".ts",".sql",".sh",".yaml",".yml",".ini",".conf",".env",".toml",".go",".rs",".java",".c",".cpp",".h",".rb",".php"}
QUVVAT_FORMATLAR = list(_FORMAT_MAP.keys()) + list(_TXT_EXTS)

def hujjat_oqi(data: bytes, fayl_nomi: str) -> dict:
    cache_tozalash()  # Eski cache larni tozalash
    fn = fayl_nomi.lower()
    for ext, (tur, func) in _FORMAT_MAP.items():
        if fn.endswith(ext):
            r = func(data)
            r["tur"] = tur
            return r
    for ext in _TXT_EXTS:
        if fn.endswith(ext):
            try:
                matn = data.decode("utf-8", errors="replace")
                r = _matn_cache(matn, "txt", data)
                r["tur"] = "txt"; r["qator_soni"] = matn.count("\n") + 1
                return r
            except: return {"tur": "txt", "xato": "O'qib bo'lmadi"}
    try:
        matn = data.decode("utf-8", errors="replace")
        if matn.strip():
            r = _matn_cache(matn, "txt", data)
            r["tur"] = "txt"; return r
    except: pass
    return {"tur": "noaniq", "xato": f"'{fayl_nomi}' qo'llab-quvvatlanmaydi.\nFormatlar: PDF, Word, Excel, EPUB, PowerPoint, FB2, RTF, HTML, JSON, MD, kod fayllar"}


# ═══════════════════════════════════════════════════════════════
#  SAHIFA O'QISH — DISKDAN
# ═══════════════════════════════════════════════════════════════

def _sahifa_oqi(h: dict, sahifa: int) -> str:
    cache = h.get("_cache")
    if cache: return cache.sahifa_oqi(sahifa)
    return h.get("sahifalar", {}).get(sahifa, "")

def sahifa_matn(h: dict, sahifa_raqam: int) -> str:
    jami = h.get("sahifalar_soni", 0)
    if sahifa_raqam < 1 or sahifa_raqam > jami:
        return f"❌ Sahifa {sahifa_raqam} topilmadi. Jami {jami} sahifa bor."
    matn = _sahifa_oqi(h, sahifa_raqam)
    if not matn: return f"📄 {sahifa_raqam}-sahifa bo'sh."
    if len(matn) > 3800: matn = matn[:3800] + "\n\n_...(qisqartirildi)_"
    return f"📄 *{sahifa_raqam}-SAHIFA* (jami {jami}):\n\n{matn}"


# ═══════════════════════════════════════════════════════════════
#  IZLASH — SQL ORQALI (100K da ham tez)
# ═══════════════════════════════════════════════════════════════

def hujjatdan_izlash(h: dict, sorov: str) -> str:
    s = sorov.lower().strip()
    m = re.search(r'(\d+)\s*[-\s]?\s*(bet|sahifa|page|стр)', s)
    if m: return sahifa_matn(h, int(m.group(1)))
    if re.match(r'^\d{1,6}$', s.strip()): return sahifa_matn(h, int(s.strip()))
    
    if any(k in s for k in ("mundarija","tarkib","содержание","bob","chapter")):
        mun = h.get("mundarija", [])
        if mun:
            t = f"📋 *MUNDARIJA ({len(mun)} bo'lim):*\n\n"
            for m in mun[:25]: t += f"  📌 {m.get('sahifa','')}-bet: {m['sarlavha']}\n"
            if len(mun) > 25: t += f"\n...va yana {len(mun)-25} ta"
            return t.rstrip()
        return "📋 Mundarija topilmadi."
    
    kalit = [k for k in s.split() if len(k) > 2]
    cache = h.get("_cache")
    
    if cache:
        topilgan = cache.izla(kalit, limit=20)
    else:
        topilgan = []
        indeks = h.get("indeks", {})
        for sn, idx in indeks.items():
            if any(k in idx.get("boshi","").lower() or k in " ".join(idx.get("sozlar",[])) for k in kalit):
                topilgan.append(sn)
    
    if topilgan:
        t = f"🔍 *\"{sorov}\"* — {len(topilgan)} sahifada topildi:\n\n"
        for sn in topilgan[:8]:
            matn = _sahifa_oqi(h, sn)
            kontekst = ""
            if matn and kalit:
                idx = matn.lower().find(kalit[0])
                if idx >= 0: kontekst = matn[max(0,idx-80):min(len(matn),idx+150)].replace("\n"," ").strip()
            if kontekst: t += f"📄 *{sn}-bet:* _{kontekst}_\n\n"
            else:
                boshi = cache.indeks_oqi(sn).get("boshi","") if cache else ""
                t += f"📄 *{sn}-bet:* _{boshi}_\n\n"
        if len(topilgan) > 8: t += f"...va yana {len(topilgan)-8} sahifada"
        return t.rstrip()
    
    return f"🔍 \"{sorov}\" topilmadi. Boshqacha so'zlar bilan urinib ko'ring."


# ═══════════════════════════════════════════════════════════════
#  FORMATLASH
# ═══════════════════════════════════════════════════════════════

def _pul(v):
    try: return f"{float(v):,.0f}"
    except: return "0"

def hujjat_xulosa_matn(h: dict, fayl_nomi: str) -> str:
    tur = h.get("tur","?"); sahifalar = h.get("sahifalar_soni", 0)
    t = f"📂 *{fayl_nomi}*\n━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
    if h.get("xato"): return t + f"❌ {h['xato']}"
    TUR_NOM = {"pdf":"PDF","docx":"Word","xlsx":"Excel","pptx":"PowerPoint","epub":"EPUB kitob",
               "fb2":"FB2 kitob","rtf":"RTF","html":"HTML","md":"Markdown","json":"JSON","odt":"ODT","txt":"Matn"}
    TUR_EMOJI = {"pdf":"📄","docx":"📝","xlsx":"📊","pptx":"📊","epub":"📚","fb2":"📚",
                 "rtf":"📄","html":"🌐","md":"📝","json":"🔧","odt":"📄","txt":"📝"}
    t += f"{TUR_EMOJI.get(tur,'📄')} Tur: *{TUR_NOM.get(tur, tur.upper())}* | Sahifalar: *{sahifalar:,}*\n"
    if tur == "epub":
        meta = h.get("metadata",{})
        if meta.get("sarlavha"): t += f"📖 Nomi: *{meta['sarlavha']}*\n"
        if meta.get("muallif"): t += f"✍️ Muallif: {meta['muallif']}\n"
    if h.get("jadvallar"): t += f"📊 Jadvallar: *{len(h['jadvallar'])}* ta\n"
    if tur == "xlsx":
        for sh in h.get("sheetlar",[])[:5]:
            t += f"\n📋 *{sh['nom']}*: {sh['qator_soni']:,} qator × {sh['ustun_soni']} ustun\n"
            stat = sh.get("statistika",{})
            if stat.get("jami"): t += f"   Jami: {_pul(stat['jami'])} | O'rtacha: {_pul(stat['ortacha'])}\n"
    mundarija = h.get("mundarija",[])
    if mundarija:
        t += f"\n📋 *MUNDARIJA ({len(mundarija)} bo'lim):*\n"
        for m in mundarija[:12]: t += f"  📌 {m.get('sahifa','')}-bet: {m['sarlavha']}\n"
        if len(mundarija) > 12: t += f"  ...va yana {len(mundarija)-12} ta\n"
    elif sahifalar > 0:
        cache = h.get("_cache")
        if cache:
            t += "\n📋 *DASTLABKI SAHIFALAR:*\n"
            for sn in range(1, min(6, sahifalar+1)):
                idx = cache.indeks_oqi(sn)
                if idx.get("boshi"): t += f"  {sn}-bet: _{idx['boshi'][:60]}_...\n"
    t += "\n💡 *BUYRUQLAR:*\n"
    if sahifalar > 0: t += f"  \"*5*\" → 5-sahifa | \"*{sahifalar:,}*\" → oxirgi\n"
    t += "  \"*so'z*\" → izlash | \"*tushuntir*\" → AI | \"*mundarija*\" → tarkib\n"
    return t.rstrip()


# ═══════════════════════════════════════════════════════════════
#  AI SAVOL-JAVOB
# ═══════════════════════════════════════════════════════════════

async def ai_hujjat_savol(h: dict, savol: str, gemini_key: str = "") -> str:
    import asyncio
    if not gemini_key: gemini_key = os.environ.get("GEMINI_API_KEY","")
    if not gemini_key: return hujjatdan_izlash(h, savol)
    kalit = [k for k in savol.lower().split() if len(k) > 2]
    cache = h.get("_cache")
    tegishli = cache.izla(kalit, limit=15) if cache else list(range(1, min(11, h.get("sahifalar_soni",0)+1)))
    if not tegishli: tegishli = list(range(1, min(11, h.get("sahifalar_soni",0)+1)))
    parts = []; jami = 0
    for sn in sorted(tegishli):
        matn = _sahifa_oqi(h, sn)
        if matn:
            qism = f"[{sn}-SAHIFA]\n{matn[:2000]}"
            if jami + len(qism) > 25000: break
            parts.append(qism); jami += len(qism)
    kontekst = "\n\n".join(parts)
    prompt = f"Sen hujjat tahlil EKSPERTIZAN. Savolga ANIQ javob ber. Sahifa raqamini ayt. To'qib chiqarma. O'zbek tilida.\n\nHUJJAT:\n{kontekst}\n\nSAVOL: {savol}\n\nJAVOB:"
    try:
        from google import genai; from google.genai import types
        client = genai.Client(api_key=gemini_key)
        loop = asyncio.get_event_loop()
        response = await asyncio.wait_for(loop.run_in_executor(None, lambda: client.models.generate_content(
            model="gemini-2.5-pro", contents=prompt,
            config=types.GenerateContentConfig(temperature=0.2, max_output_tokens=2000))), timeout=30)
        javob = (response.text or "").strip()
        if javob: return f"🤖 *AI JAVOB:*\n\n{javob}"
    except Exception as e: log.debug("AI: %s", e)
    return hujjatdan_izlash(h, savol)


# ═══════════════════════════════════════════════════════════════
#  SAVOL TURI ANIQLASH
# ═══════════════════════════════════════════════════════════════

QIDIRUV_SOZLAR = ("qayerda","topib ber","izla","bor","найди","где","nima yozilgan","nima bor",
    "haqida","aytib ber","ko'rsat","ber","chiqar","o'qi","покажи","qancha","nechta",
    "nechanchi","mundarija","sarlavha","jadval","bob","mavzu","formula","tarkib","qism","bo'lim")

def hujjat_sorov_bormi(matn: str) -> bool:
    m = matn.lower().strip()
    if re.search(r'\d+\s*[-\s]?\s*(bet|sahifa|page|стр)', m): return True
    if re.match(r'^\d{1,6}$', m.strip()): return True
    return any(s in m for s in QIDIRUV_SOZLAR)

def ai_savol_kerakmi(matn: str) -> bool:
    m = matn.lower().strip()
    return any(s in m for s in ("tushuntir","nimaga","qanday","nima uchun","farqi","asosiy","muhim",
        "xulosa","qisqacha","formala","misol","izohla","yechim","javob","объясни","почему",
        "tahlil","solishtir","sabab","natija"))
